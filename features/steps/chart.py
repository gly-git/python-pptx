# encoding: utf-8

"""
Gherkin step implementations for chart features.
"""

from __future__ import absolute_import, print_function

import hashlib

from itertools import islice

from behave import given, then, when

from pptx import Presentation
from pptx.chart.chart import Legend
from pptx.chart.data import BubbleChartData, ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
)
from pptx.enum.dml import MSO_FILL_TYPE, MSO_THEME_COLOR
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from pptx.text.text import Font
from pptx.util import Inches

from helpers import count, test_pptx


# given ===================================================

@given('a bar plot having known categories')
def given_a_bar_plot_having_known_categories(context):
    prs = Presentation(test_pptx('cht-plot-props'))
    context.plot = prs.slides[2].shapes[0].chart.plots[0]


@given('a bar plot {having_or_not} data labels')
def given_a_bar_plot_having_or_not_data_labels(context, having_or_not):
    slide_idx = {
        'having':     0,
        'not having': 1,
    }[having_or_not]
    prs = Presentation(test_pptx('cht-plot-props'))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given('a bar plot having gap width of {width}')
def given_a_bar_plot_having_gap_width_of_width(context, width):
    slide_idx = {'no explicit value': 0, '300': 1}[width]
    prs = Presentation(test_pptx('cht-plot-props'))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given('a bar plot having overlap of {overlap}')
def given_a_bar_plot_having_overlap_of_overlap(context, overlap):
    slide_idx = {
        'no explicit value': 0,
        '42':                1,
        '-42':               2,
    }[overlap]
    prs = Presentation(test_pptx('cht-plot-props'))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given('a bar plot having vary color by category set to {setting}')
def given_a_bar_plot_having_vary_color_by_category_setting(context, setting):
    slide_idx = {
        'no explicit setting': 0,
        'True':                1,
        'False':               2,
    }[setting]
    prs = Presentation(test_pptx('cht-plot-props'))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given('a bar series having fill of {fill}')
def given_a_bar_series_having_fill_of_fill(context, fill):
    series_idx = {
        'Automatic': 0,
        'No Fill':   1,
        'Orange':    2,
        'Accent 1':  3,
    }[fill]
    prs = Presentation(test_pptx('cht-series-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series = plot.series[series_idx]


@given('a bar series having invert_if_negative of {setting}')
def given_a_bar_series_having_invert_if_negative_setting(context, setting):
    series_idx = {
        'no explicit setting': 0,
        'True':                1,
        'False':               2,
    }[setting]
    prs = Presentation(test_pptx('cht-series-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series = plot.series[series_idx]


@given('a bar series having known values')
def given_a_bar_series_having_known_values(context):
    prs = Presentation(test_pptx('cht-series-props'))
    context.series = prs.slides[1].shapes[0].chart.plots[0].series[0]


@given('a bar series having {width} line')
def given_a_bar_series_having_width_line(context, width):
    series_idx = {
        'no':      0,
        '1 point': 1,
    }[width]
    prs = Presentation(test_pptx('cht-series-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series = plot.series[series_idx]


@given('a chart {having_or_not} a legend')
def given_a_chart_having_or_not_a_legend(context, having_or_not):
    slide_idx = {
        'having':     0,
        'not having': 1,
    }[having_or_not]
    prs = Presentation(test_pptx('cht-legend'))
    context.chart = prs.slides[slide_idx].shapes[0].chart


@given('a chart of size and type {spec}')
def given_a_chart_of_size_and_type_spec(context, spec):
    slide_idx = {
        '2x2 Clustered Bar':    0,
        '2x2 100% Stacked Bar': 1,
        '2x2 Clustered Column': 2,
        '4x3 Line':             3,
        '3x1 Pie':              4,
    }[spec]
    prs = Presentation(test_pptx('cht-replace-data'))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.chart = chart
    context.xlsx_sha1 = hashlib.sha1(
        chart._workbook.xlsx_part.blob
    ).hexdigest()


@given('a chart of type {chart_type}')
def given_a_chart_of_type_chart_type(context, chart_type):
    slide_idx, shape_idx = {
        'Area':                     (0, 0),
        'Stacked Area':             (0, 1),
        '100% Stacked Area':        (0, 2),
        '3-D Area':                 (0, 3),
        '3-D Stacked Area':         (0, 4),
        '3-D 100% Stacked Area':    (0, 5),
        'Clustered Bar':            (1, 0),
        'Stacked Bar':              (1, 1),
        '100% Stacked Bar':         (1, 2),
        'Clustered Column':         (1, 3),
        'Stacked Column':           (1, 4),
        '100% Stacked Column':      (1, 5),
        'Line':                     (2, 0),
        'Stacked Line':             (2, 1),
        '100% Stacked Line':        (2, 2),
        'Marked Line':              (2, 3),
        'Stacked Marked Line':      (2, 4),
        '100% Stacked Marked Line': (2, 5),
        'Pie':                      (3, 0),
        'Exploded Pie':             (3, 1),
        'XY (Scatter)':             (4, 0),
        'XY Lines':                 (4, 1),
        'XY Lines No Markers':      (4, 2),
        'XY Smooth Lines':          (4, 3),
        'XY Smooth No Markers':     (4, 4),
        'Bubble':                   (5, 0),
        '3D-Bubble':                (5, 1),
    }[chart_type]
    prs = Presentation(test_pptx('cht-chart-type'))
    context.chart = prs.slides[slide_idx].shapes[shape_idx].chart


@given('a data label {having_or_not} custom text')
def given_a_data_label_having_or_not_custom_text(context, having_or_not):
    point_idx = {
        'having':    0,
        'having no': 1,
    }[having_or_not]
    prs = Presentation(test_pptx('cht-point-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.data_label = plot.series[0].points[point_idx].data_label


@given('a data label positioned {relation} its data point')
def given_a_data_label_positioned_relation_its_data_point(context, relation):
    point_idx = {
        'in unspecified relation to': 0,
        'centered on':                1,
        'below':                      2,
    }[relation]
    prs = Presentation(test_pptx('cht-point-props'))
    plot = prs.slides[1].shapes[0].chart.plots[0]
    context.data_label = plot.series[0].points[point_idx].data_label


@given('a legend')
def given_a_legend(context):
    prs = Presentation(test_pptx('cht-legend-props'))
    context.legend = prs.slides[0].shapes[0].chart.legend


@given('a legend having horizontal offset of {value}')
def given_a_legend_having_horizontal_offset_of_value(context, value):
    slide_idx = {
        'none': 0,
        '-0.5': 1,
        '0.42': 2,
    }[value]
    prs = Presentation(test_pptx('cht-legend-props'))
    context.legend = prs.slides[slide_idx].shapes[0].chart.legend


@given('a legend positioned {location} the chart')
def given_a_legend_positioned_location_the_chart(context, location):
    slide_idx = {
        'at an unspecified location of': 0,
        'below':                         1,
        'to the right of':               2,
    }[location]
    prs = Presentation(test_pptx('cht-legend-props'))
    context.legend = prs.slides[slide_idx].shapes[0].chart.legend


@given('a legend with overlay setting of {setting}')
def given_a_legend_with_overlay_setting_of_setting(context, setting):
    slide_idx = {
        'no explicit setting': 0,
        'True':                1,
        'False':               2,
    }[setting]
    prs = Presentation(test_pptx('cht-legend-props'))
    context.legend = prs.slides[slide_idx].shapes[0].chart.legend


@given('a point')
def given_a_point(context):
    prs = Presentation(test_pptx('cht-point-props'))
    chart = prs.slides[0].shapes[0].chart
    context.point = chart.plots[0].series[0].points[0]


@given('a {points_type} object containing 3 points')
def given_a_points_type_object_containing_3_points(context, points_type):
    slide_idx = {'XyPoints': 0, 'BubblePoints': 1}[points_type]
    prs = Presentation(test_pptx('cht-point-access'))
    series = prs.slides[slide_idx].shapes[0].chart.plots[0].series[0]
    context.points = series.points


@given('a series of type {series_type}')
def given_a_series_of_type_series_type(context, series_type):
    slide_idx = {'XY': 2, 'Bubble': 3}[series_type]
    prs = Presentation(test_pptx('cht-series-props'))
    context.series = prs.slides[slide_idx].shapes[0].chart.plots[0].series[0]


@given('an axis')
def given_an_axis(context):
    prs = Presentation(test_pptx('cht-axis-props'))
    chart = prs.slides[0].shapes[0].chart
    context.axis = chart.value_axis


@given('an axis having {major_or_minor} gridlines')
def given_an_axis_having_major_or_minor_gridlines(context, major_or_minor):
    prs = Presentation(test_pptx('cht-axis-props'))
    chart = prs.slides[0].shapes[0].chart
    context.axis = chart.value_axis


@given('an axis having {major_or_minor} unit of {value}')
def given_an_axis_having_major_or_minor_unit_of_value(
        context, major_or_minor, value):
    slide_idx = 0 if value == 'Auto' else 1
    prs = Presentation(test_pptx('cht-axis-props'))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.axis = chart.value_axis


@given('an axis not having {major_or_minor} gridlines')
def given_an_axis_not_having_major_or_minor_gridlines(context, major_or_minor):
    prs = Presentation(test_pptx('cht-axis-props'))
    chart = prs.slides[0].shapes[0].chart
    context.axis = chart.category_axis


@given('bar chart data labels positioned {relation_to} their data point')
def given_bar_chart_data_labels_positioned_relation_to_their_data_point(
        context, relation_to):
    slide_idx = {
        'in unspecified relation to': 0,
        'inside, at the base of':     1,
    }[relation_to]
    prs = Presentation(test_pptx('cht-datalabels-props'))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.data_labels = chart.plots[0].data_labels


@given('tick labels having an offset of {setting}')
def given_tick_labels_having_an_offset_of_setting(context, setting):
    slide_idx = {
        'no explicit setting': 0,
        '420':                 1,
    }[setting]
    prs = Presentation(test_pptx('cht-ticklabels-props'))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.tick_labels = chart.category_axis.tick_labels


# when ====================================================

@when('I add a {kind} chart with {cats} categories and {sers} series')
def when_I_add_a_chart_with_categories_and_series(context, kind, cats, sers):
    chart_type = {
        'Clustered Bar':    XL_CHART_TYPE.BAR_CLUSTERED,
        '100% Stacked Bar': XL_CHART_TYPE.BAR_STACKED_100,
        'Clustered Column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'Line':             XL_CHART_TYPE.LINE,
        'Pie':              XL_CHART_TYPE.PIE,
    }[kind]
    category_count, series_count = int(cats), int(sers)
    category_source = ('Foo', 'Bar', 'Baz', 'Boo', 'Far', 'Faz')
    series_value_source = count(1.1, 1.1)

    chart_data = ChartData()
    chart_data.categories = category_source[:category_count]
    for idx in range(series_count):
        series_title = 'Series %d' % (idx+1)
        series_values = tuple(islice(series_value_source, category_count))
        chart_data.add_series(series_title, series_values)

    context.chart = context.slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1), Inches(8), Inches(5), chart_data
    ).chart


@when('I add a {bubble_type} chart having 2 series of 3 points each')
def when_I_add_a_bubble_chart_having_2_series_of_3_pts(context, bubble_type):
    chart_type = getattr(XL_CHART_TYPE, bubble_type)
    data = (
        ('Series 1', ((-0.1, 0.5, 1.0), (16.2, 0.0, 2.0), (8.0, -0.2, 3.0))),
        ('Series 2', ((12.4, 0.8, 4.0), (-7.5, 0.5, 5.0), (5.1, -0.5, 6.0))),
    )

    chart_data = BubbleChartData()

    for series_data in data:
        series_label, points = series_data
        series = chart_data.add_series(series_label)
        for point in points:
            x, y, size = point
            series.add_data_point(x, y, size)

    context.chart = context.slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1), Inches(8), Inches(5), chart_data
    ).chart


@when('I add an {xy_type} chart having 2 series of 3 points each')
def when_I_add_an_xy_chart_having_2_series_of_3_points(context, xy_type):
    chart_type = getattr(XL_CHART_TYPE, xy_type)
    data = (
        ('Series 1', ((-0.1, 0.5), (16.2, 0.0), (8.0,  0.2))),
        ('Series 2', ((12.4, 0.8), (-7.5, -0.5), (-5.1, -0.2)))
    )

    chart_data = XyChartData()

    for series_data in data:
        series_label, points = series_data
        series = chart_data.add_series(series_label)
        for point in points:
            x, y = point
            series.add_data_point(x, y)

    context.chart = context.slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1), Inches(8), Inches(5), chart_data
    ).chart


@when('I assign {value} to chart.has_legend')
def when_I_assign_value_to_chart_has_legend(context, value):
    new_value = {
        'True':  True,
        'False': False,
    }[value]
    context.chart.has_legend = new_value


@when('I assign {value} to axis.has_{major_or_minor}_gridlines')
def when_I_assign_value_to_axis_has_major_or_minor_gridlines(
        context, value, major_or_minor):
    axis = context.axis
    propname = 'has_%s_gridlines' % major_or_minor
    new_value = {'True': True, 'False': False}[value]
    setattr(axis, propname, new_value)


@when('I assign {value} to axis.{major_or_minor}_unit')
def when_I_assign_value_to_axis_major_or_minor_unit(
        context, value, major_or_minor):
    axis = context.axis
    propname = '%s_unit' % major_or_minor
    new_value = {'8.4': 8.4, '5': 5, 'None': None}[value]
    setattr(axis, propname, new_value)


@when('I assign {value} to data_label.has_text_frame')
def when_I_assign_value_to_data_label_has_text_frame(context, value):
    new_value = {'True': True, 'False': False}[value]
    context.data_label.has_text_frame = new_value


@when('I assign {value} to data_label.position')
def when_I_assign_value_to_data_label_position(context, value):
    new_value = (
        None if value == 'None' else getattr(XL_DATA_LABEL_POSITION, value)
    )
    context.data_label.position = new_value


@when('I assign {value} to data_labels.position')
def when_I_assign_value_to_data_labels_position(context, value):
    new_value = (
        None if value == 'None' else getattr(XL_DATA_LABEL_POSITION, value)
    )
    context.data_labels.position = new_value


@when('I assign {value} to legend.horz_offset')
def when_I_assign_value_to_legend_horz_offset(context, value):
    new_value = float(value)
    context.legend.horz_offset = new_value


@when('I assign {value} to legend.include_in_layout')
def when_I_assign_value_to_legend_include_in_layout(context, value):
    new_value = {
        'True':  True,
        'False': False,
    }[value]
    context.legend.include_in_layout = new_value


@when('I assign {value} to legend.position')
def when_I_assign_value_to_legend_position(context, value):
    enum_value = getattr(XL_LEGEND_POSITION, value)
    context.legend.position = enum_value


@when('I assign {value} to plot.gap_width')
def when_I_assign_value_to_plot_gap_width(context, value):
    new_value = int(value)
    context.plot.gap_width = new_value


@when('I assign {value} to plot.has_data_labels')
def when_I_assign_value_to_plot_has_data_labels(context, value):
    new_value = {
        'True':  True,
        'False': False,
    }[value]
    context.plot.has_data_labels = new_value


@when('I assign {value} to plot.overlap')
def when_I_assign_value_to_plot_overlap(context, value):
    new_value = int(value)
    context.plot.overlap = new_value


@when('I assign {value} to plot.vary_by_categories')
def when_I_assign_value_to_plot_vary_by_categories(context, value):
    new_value = {
        'True':  True,
        'False': False,
    }[value]
    context.plot.vary_by_categories = new_value


@when('I assign {value} to series.invert_if_negative')
def when_I_assign_value_to_series_invert_if_negative(context, value):
    new_value = {
        'True':  True,
        'False': False,
    }[value]
    context.series.invert_if_negative = new_value


@when('I assign {value} to tick_labels.offset')
def when_I_assign_value_to_tick_labels_offset(context, value):
    new_value = int(value)
    context.tick_labels.offset = new_value


@when('I replace its data with {cats} categories and {sers} series')
def when_I_replace_its_data_with_categories_and_series(context, cats, sers):
    category_count, series_count = int(cats), int(sers)
    category_source = ('Foo', 'Bar', 'Baz', 'Boo', 'Far', 'Faz')
    series_value_source = count(1.1, 1.1)

    chart_data = ChartData()
    chart_data.categories = category_source[:category_count]
    for idx in range(series_count):
        series_title = 'New Series %d' % (idx+1)
        series_values = tuple(islice(series_value_source, category_count))
        chart_data.add_series(series_title, series_values)

    context.chart.replace_data(chart_data)


# then ====================================================

@then('axis.has_{major_or_minor}_gridlines is {value}')
def then_axis_has_major_or_minor_gridlines_is_expected_value(
        context, major_or_minor, value):
    axis = context.axis
    actual_value = {
        'major': axis.has_major_gridlines,
        'minor': axis.has_minor_gridlines,
    }[major_or_minor]
    expected_value = {'True': True, 'False': False}[value]
    assert actual_value is expected_value, 'got %s' % actual_value


@then('axis.major_gridlines is a MajorGridlines object')
def then_axis_major_gridlines_is_a_MajorGridlines_object(context):
    axis = context.axis
    assert type(axis.major_gridlines).__name__ == 'MajorGridlines'


@then('axis.{major_or_minor}_unit is {value}')
def then_axis_major_or_minor_unit_is_value(context, major_or_minor, value):
    axis = context.axis
    propname = '%s_unit' % major_or_minor
    actual_value = getattr(axis, propname)
    expected_value = {
        '20.0': 20.0, '8.4': 8.4, '5.0': 5.0, '4.2': 4.2, 'None': None
    }[value]
    assert actual_value == expected_value, 'got %s' % actual_value


@then('chart.category_axis is a {type_name} object')
def then_chart_category_axis_is_an_axis_type_object(context, type_name):
    category_axis = context.chart.category_axis
    assert type(category_axis).__name__ == type_name


@then('chart.chart_type is {enum_member}')
def then_chart_chart_type_is_value(context, enum_member):
    expected_value = getattr(XL_CHART_TYPE, enum_member)
    chart = context.chart
    assert chart.chart_type is expected_value, 'got %s' % chart.chart_type


@then('chart.has_legend is {value}')
def then_chart_has_legend_is_value(context, value):
    expected_value = {
        'True':  True,
        'False': False,
    }[value]
    chart = context.chart
    assert chart.has_legend is expected_value


@then('chart.legend is a legend object')
def then_chart_legend_is_a_legend_object(context):
    chart = context.chart
    assert isinstance(chart.legend, Legend)


@then('chart.value_axis is a ValueAxis object')
def then_chart_value_axis_is_a_ValueAxis_object(context):
    value_axis = context.chart.value_axis
    assert type(value_axis).__name__ == 'ValueAxis'


@then('data_label.has_text_frame is {value}')
def then_data_label_has_text_frame_is_value(context, value):
    expected_value = {'True': True, 'False': False}[value]
    data_label = context.data_label
    assert data_label.has_text_frame is expected_value


@then('data_label.position is {value}')
def then_data_label_position_is_value(context, value):
    expected_value = (
        None if value == 'None' else getattr(XL_DATA_LABEL_POSITION, value)
    )
    data_label = context.data_label
    assert data_label.position is expected_value, (
        'got %s' % data_label.position
    )


@then('data_label.text_frame is a TextFrame object')
def then_data_label_text_frame_is_a_TextFrame_object(context):
    text_frame = context.data_label.text_frame
    assert type(text_frame).__name__ == 'TextFrame'


@then('data_labels.position is {value}')
def then_data_labels_position_is_value(context, value):
    expected_value = (
        None if value == 'None' else getattr(XL_DATA_LABEL_POSITION, value)
    )
    data_labels = context.data_labels
    assert data_labels.position is expected_value, (
        'got %s' % data_labels.position
    )


@then('each series has a new name')
def then_each_series_has_a_new_name(context):
    for series in context.chart.plots[0].series:
        assert series.name.startswith('New ')


@then('each series has {count} values')
def then_each_series_has_count_values(context, count):
    expected_count = int(count)
    for series in context.chart.plots[0].series:
        actual_value_count = len(series.values)
        assert actual_value_count == expected_count


@then('iterating points produces 3 Point objects')
def then_iterating_points_produces_3_point_objects(context):
    points = context.points
    idx = -1
    for idx, point in enumerate(points):
        assert type(point).__name__ == 'Point'
    assert idx == 2


@then('legend.font is a Font object')
def then_legend_font_is_a_Font_object(context):
    legend = context.legend
    assert isinstance(legend.font, Font)


@then('legend.horz_offset is {value}')
def then_legend_horz_offset_is_value(context, value):
    expected_value = float(value)
    legend = context.legend
    assert legend.horz_offset == expected_value


@then('legend.include_in_layout is {value}')
def then_legend_include_in_layout_is_value(context, value):
    expected_value = {
        'True':  True,
        'False': False,
    }[value]
    legend = context.legend
    assert legend.include_in_layout is expected_value


@then('legend.position is {value}')
def then_legend_position_is_value(context, value):
    expected_position = getattr(XL_LEGEND_POSITION, value)
    legend = context.legend
    assert legend.position is expected_position, 'got %s' % legend.position


@then('len(chart.series) is {count}')
def then_len_chart_series_is_count(context, count):
    expected_count = int(count)
    assert len(context.chart.series) == expected_count


@then('len(plot.categories) is {count}')
def then_len_plot_categories_is_count(context, count):
    plot = context.chart.plots[0]
    expected_count = int(count)
    assert len(plot.categories) == expected_count


@then('len(points) is 3')
def then_len_points_is_3(context):
    points = context.points
    assert len(points) == 3


@then('len(series.values) is {count} for each series')
def then_len_series_values_is_count_for_each_series(context, count):
    expected_count = int(count)
    for series in context.chart.plots[0].series:
        assert len(series.values) == expected_count


@then('plot.categories contains the known category strings')
def then_plot_categories_contains_the_known_category_strings(context):
    plot = context.plot
    expected_categories = ('Foo', 'Bar', 'Baz')
    assert plot.categories == expected_categories, (
        'got %s' % plot.categories
    )


@then('plot.gap_width is {value}')
def then_plot_gap_width_is_value(context, value):
    expected_value = int(value)
    plot = context.plot
    assert plot.gap_width == expected_value, 'got %s' % plot.gap_width


@then('plot.has_data_labels is {value}')
def then_plot_has_data_labels_is_value(context, value):
    expected_value = {
        'True':  True,
        'False': False,
    }[value]
    assert context.plot.has_data_labels is expected_value


@then('plot.overlap is {value}')
def then_plot_overlap_is_expected_value(context, value):
    expected_value = int(value)
    plot = context.plot
    assert plot.overlap == expected_value, 'got %s' % plot.overlap


@then('plot.vary_by_categories is {value}')
def then_plot_vary_by_categories_is_value(context, value):
    expected_value = {
        'True':  True,
        'False': False,
    }[value]
    plot = context.plot
    assert plot.vary_by_categories is expected_value


@then('point.data_label is a DataLabel object')
def then_point_data_label_is_a_DataLabel_object(context):
    point = context.point
    assert type(point.data_label).__name__ == 'DataLabel'


@then('points[2] is a Point object')
def then_points_2_is_a_Point_object(context):
    points = context.points
    assert type(points[2]).__name__ == 'Point'


@then('series.fill.fore_color.rgb is FF6600')
def then_series_fill_fore_color_rgb_is_FF6600(context):
    fill = context.series.fill
    assert fill.fore_color.rgb == RGBColor(0xFF, 0x66, 0x00)


@then('series.fill.fore_color.theme_color is Accent 1')
def then_series_fill_fore_color_theme_color_is_Accent_1(context):
    fill = context.series.fill
    assert fill.fore_color.theme_color == MSO_THEME_COLOR.ACCENT_1


@then('series.fill.type is {fill_type}')
def then_series_fill_type_is_type(context, fill_type):
    expected_fill_type = {
        'None':                     None,
        'MSO_FILL_TYPE.BACKGROUND': MSO_FILL_TYPE.BACKGROUND,
        'MSO_FILL_TYPE.SOLID':      MSO_FILL_TYPE.SOLID,
    }[fill_type]
    fill = context.series.fill
    assert fill.type == expected_fill_type


@then('series.invert_if_negative is {value}')
def then_series_invert_if_negative_is_value(context, value):
    expected_value = {
        'True':  True,
        'False': False,
    }[value]
    series = context.series
    assert series.invert_if_negative is expected_value


@then('series.line.width is {width}')
def then_series_line_width_is_width(context, width):
    expected_width = int(width)
    line = context.series.line
    assert line.width == expected_width


@then('series.points is a {type_name} object')
def then_series_points_is_a_type_name_object(context, type_name):
    series = context.series
    assert type(series.points).__name__ == type_name


@then('series.values contains the known values')
def then_series_values_contains_the_known_values(context):
    series = context.series
    expected_values = (1.2, 2.3, 3.4)
    assert series.values == expected_values, 'got %s' % series.values


@then('the chart has an Excel data worksheet')
def then_the_chart_has_an_Excel_data_worksheet(context):
    xlsx_part = context.chart._workbook.xlsx_part
    assert isinstance(xlsx_part, EmbeddedXlsxPart)


@then('the chart has new chart data')
def then_the_chart_has_new_chart_data(context):
    orig_xlsx_sha1 = context.xlsx_sha1
    new_xlsx_sha1 = hashlib.sha1(
        context.chart._workbook.xlsx_part.blob
    ).hexdigest()
    assert new_xlsx_sha1 != orig_xlsx_sha1


@then('tick_labels.offset is {value}')
def then_tick_labels_offset_is_expected_value(context, value):
    expected_value = int(value)
    tick_labels = context.tick_labels
    assert tick_labels.offset == expected_value, (
        'got %s' % tick_labels.offset
    )
